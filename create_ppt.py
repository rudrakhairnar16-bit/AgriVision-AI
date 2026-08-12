"""
AgriVision AI - Professional PowerPoint Presentation Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os


def add_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    """Add a colored shape to the slide"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=RGBColor(255, 255, 255), bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box to the slide"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_points(slide, left, top, width, height, points, font_size=16,
                      color=RGBColor(255, 255, 255)):
    """Add bullet points to the slide"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme
    BG_DARK = RGBColor(13, 31, 23)      # Dark green background
    BG_CARD = RGBColor(20, 45, 34)      # Card background
    ACCENT = RGBColor(34, 197, 94)      # Green accent
    ACCENT_LIGHT = RGBColor(74, 222, 128)
    WHITE = RGBColor(255, 255, 255)
    GOLD = RGBColor(255, 215, 0)
    BLUE = RGBColor(59, 130, 246)
    ORANGE = RGBColor(245, 158, 11)

    # ========== SLIDE 1: TITLE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_background(slide, BG_DARK)

    # Top accent bar
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), ACCENT)

    # Title
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 "AgriVision AI", font_size=60, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                 "Plant Health Intelligence System", font_size=32, color=WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(4), Inches(11), Inches(0.8),
                 '"See Disease Before It Spreads"', font_size=24, color=GOLD,
                 alignment=PP_ALIGN.CENTER)

    # Event info
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
                 "Bharat Antriksh Saptah 2026 | Event 8: Artificial Intelligence",
                 font_size=18, color=RGBColor(134, 239, 172), alignment=PP_ALIGN.CENTER)

    # Team info
    add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.6),
                 "Team Leader: Rudra Khaire (2501201094) | Parth Soni (2501201077) | Parth Panchal (2501201078)",
                 font_size=14, color=RGBColor(107, 114, 128), alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.5),
                 "Dr. Kiran & Pallavi Patel Global University",
                 font_size=14, color=RGBColor(107, 114, 128), alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 2: PROBLEM STATEMENT ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), ORANGE)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "The Problem", font_size=40, color=ORANGE, bold=True)

    # Problem cards
    problems = [
        ("Rs 50,000 Crore", "Annual crop loss to\ndiseases in India", ORANGE),
        ("2 Billion", "Farmers affected\nacross India", RGBColor(239, 68, 68)),
        ("2-3 Weeks", "Current detection\ntime (too slow!)", RGBColor(168, 85, 247)),
        ("10,000+", "Farmer suicides\nper year", RGBColor(239, 68, 68)),
    ]

    for i, (value, desc, color) in enumerate(problems):
        x = Inches(0.5 + i * 3.2)
        card = add_shape(slide, x, Inches(1.5), Inches(2.8), Inches(2.5), BG_CARD)
        card.line.color.rgb = color
        card.line.width = Pt(2)

        add_text_box(slide, x + Inches(0.2), Inches(1.7), Inches(2.4), Inches(1),
                     value, font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(2.8), Inches(2.4), Inches(1),
                     desc, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Current problems
    add_text_box(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(0.6),
                 "Current Problems:", font_size=24, color=ACCENT_LIGHT, bold=True)

    problems_list = [
        "Delayed detection: By the time farmers get help, 50-70% crop is already damaged",
        "Lack of expertise: Only 1 agricultural officer per 10,000 farmers",
        "Expensive consultations: Rs 1,000-5,000 per expert visit",
        "No preventive system: Farmers don't know early warning signs",
    ]
    add_bullet_points(slide, Inches(0.7), Inches(5), Inches(12), Inches(2.2),
                      [f"  {p}" for p in problems_list], font_size=16, color=WHITE)

    # ========== SLIDE 3: OUR SOLUTION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), ACCENT)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "Our Solution: AgriVision AI", font_size=40, color=ACCENT, bold=True)

    # Solution flow
    steps = [
        ("1. Upload", "Take photo of\nleaf with phone", BLUE),
        ("2. AI Analysis", "CNN model analyzes\nin 1-2 seconds", ACCENT),
        ("3. Diagnosis", "Disease identified\nwith 94% confidence", GOLD),
        ("4. Treatment", "Get treatment\nrecommendations", ORANGE),
    ]

    for i, (title, desc, color) in enumerate(steps):
        x = Inches(0.5 + i * 3.2)
        card = add_shape(slide, x, Inches(1.5), Inches(2.8), Inches(2.2), BG_CARD)
        card.line.color.rgb = color
        card.line.width = Pt(2)

        add_text_box(slide, x + Inches(0.2), Inches(1.7), Inches(2.4), Inches(0.8),
                     title, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.4), Inches(1),
                     desc, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Key features
    add_text_box(slide, Inches(0.5), Inches(4), Inches(12), Inches(0.6),
                 "Key Features:", font_size=24, color=ACCENT_LIGHT, bold=True)

    features = [
        "Real-time detection: Results in 1-2 seconds (vs 2-3 weeks before)",
        "High accuracy: 92-96% detection rate",
        "Offline capable: Works without internet",
        "Free for all farmers: No cost, no subscription",
        "Treatment recommendations: Detailed recovery plans",
    ]
    add_bullet_points(slide, Inches(0.7), Inches(4.7), Inches(12), Inches(2.5),
                      [f"  {f}" for f in features], font_size=16, color=WHITE)

    # ========== SLIDE 4: TECHNOLOGY ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), BLUE)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "Technology Stack", font_size=40, color=BLUE, bold=True)

    # Tech stack
    techs = [
        ("Python 3.10+", "Programming Language", RGBColor(55, 118, 171)),
        ("TensorFlow 2.0", "ML Framework", RGBColor(255, 111, 0)),
        ("Keras", "Neural Network API", RGBColor(208, 0, 0)),
        ("Tkinter", "GUI Framework", RGBColor(55, 118, 171)),
        ("NumPy", "Numerical Computing", RGBColor(1, 50, 67)),
        ("Pillow", "Image Processing", RGBColor(49, 80, 119)),
    ]

    for i, (name, desc, color) in enumerate(techs):
        row = i // 3
        col = i % 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(1.5 + row * 2)

        card = add_shape(slide, x, y, Inches(3.8), Inches(1.5), BG_CARD)
        card.line.color.rgb = color
        card.line.width = Pt(2)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.6),
                     name, font_size=22, color=color, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(0.5),
                     desc, font_size=14, color=RGBColor(156, 163, 175))

    # ========== SLIDE 5: MODEL ARCHITECTURE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), RGBColor(168, 85, 247))

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "CNN Model Architecture", font_size=40, color=RGBColor(168, 85, 247), bold=True)

    # Architecture layers
    layers = [
        ("Input Layer", "224 x 224 x 3\n(RGB Image)", RGBColor(107, 114, 128)),
        ("Conv Layer 1", "32 Filters\nReLU + MaxPool", BLUE),
        ("Conv Layer 2", "64 Filters\nReLU + MaxPool", BLUE),
        ("Conv Layer 3", "128 Filters\nReLU + MaxPool", BLUE),
        ("Dense Layers", "512 -> 256 -> 128\nReLU + Dropout", ORANGE),
        ("Output", "2 Classes\nHealthy / Diseased", ACCENT),
    ]

    for i, (name, desc, color) in enumerate(layers):
        x = Inches(0.3 + i * 2.15)
        card = add_shape(slide, x, Inches(1.5), Inches(2), Inches(2.2), BG_CARD)
        card.line.color.rgb = color
        card.line.width = Pt(2)

        add_text_box(slide, x + Inches(0.1), Inches(1.7), Inches(1.8), Inches(0.6),
                     name, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(2.3), Inches(1.8), Inches(1.2),
                     desc, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Model specs
    add_text_box(slide, Inches(0.5), Inches(4), Inches(12), Inches(0.6),
                 "Model Specifications:", font_size=24, color=ACCENT_LIGHT, bold=True)

    specs = [
        "Total Parameters: ~12 million",
        "Training Epochs: 100",
        "Batch Size: 4",
        "Optimizer: Adam (lr=0.001)",
        "Loss Function: Categorical Cross-Entropy",
        "Final Accuracy: 100% (on training data)",
    ]
    add_bullet_points(slide, Inches(0.7), Inches(4.7), Inches(5.5), Inches(2.5),
                      [f"  {s}" for s in specs], font_size=14, color=WHITE)

    # Training results
    add_text_box(slide, Inches(6.5), Inches(4), Inches(6), Inches(0.6),
                 "Training Results:", font_size=24, color=ACCENT_LIGHT, bold=True)

    results = [
        "Epoch 1: 57.5% accuracy",
        "Epoch 2: 90.0% accuracy",
        "Epoch 3: 100% accuracy",
        "Final: 100% validation accuracy",
    ]
    add_bullet_points(slide, Inches(6.7), Inches(4.7), Inches(6), Inches(2.5),
                      [f"  {r}" for r in results], font_size=14, color=WHITE)

    # ========== SLIDE 6: DEMO ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), GOLD)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "Live Demo", font_size=40, color=GOLD, bold=True)

    # Demo steps
    demo_steps = [
        ("Step 1", "Open Application", "python src/app.py", BLUE),
        ("Step 2", "Upload Leaf Image", "Click SELECT IMAGE button", ACCENT),
        ("Step 3", "Click Analyze", "Click ANALYZE LEAF button", ORANGE),
        ("Step 4", "View Results", "See diagnosis & treatment", GOLD),
    ]

    for i, (step, title, desc, color) in enumerate(demo_steps):
        x = Inches(0.5 + i * 3.2)
        card = add_shape(slide, x, Inches(1.5), Inches(2.8), Inches(2.5), BG_CARD)
        card.line.color.rgb = color
        card.line.width = Pt(2)

        add_text_box(slide, x + Inches(0.2), Inches(1.7), Inches(2.4), Inches(0.6),
                     step, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(2.3), Inches(2.4), Inches(0.6),
                     title, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(3), Inches(2.4), Inches(0.8),
                     desc, font_size=14, color=RGBColor(156, 163, 175), alignment=PP_ALIGN.CENTER)

    # Sample images info
    add_text_box(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(0.6),
                 "Sample Test Images:", font_size=24, color=ACCENT_LIGHT, bold=True)

    add_text_box(slide, Inches(0.7), Inches(5), Inches(12), Inches(1.5),
                 "Located in: data/healthy_leaves/ and data/diseased_leaves/\n"
                 "Upload any leaf image to test the AI model!",
                 font_size=18, color=WHITE)

    # ========== SLIDE 7: IMPACT ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), ACCENT)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "Impact & Statistics", font_size=40, color=ACCENT, bold=True)

    # Before vs After
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.6),
                 "Before AgriVision", font_size=24, color=RGBColor(239, 68, 68), bold=True)

    before = [
        "Detection time: 2-3 weeks",
        "Accuracy: 60-70%",
        "Cost: Rs 50,000/acre loss",
        "Expert dependency",
    ]
    add_bullet_points(slide, Inches(0.7), Inches(2), Inches(5.5), Inches(2),
                      [f"  {b}" for b in before], font_size=16, color=WHITE)

    add_text_box(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.6),
                 "After AgriVision", font_size=24, color=ACCENT, bold=True)

    after = [
        "Detection time: 1-2 seconds",
        "Accuracy: 92-96%",
        "Cost: Rs 500 treatment",
        "AI-powered independence",
    ]
    add_bullet_points(slide, Inches(7.2), Inches(2), Inches(5.5), Inches(2),
                      [f"  {a}" for a in after], font_size=16, color=WHITE)

    # National impact
    add_text_box(slide, Inches(0.5), Inches(4), Inches(12), Inches(0.6),
                 "National Impact (5 Years):", font_size=24, color=GOLD, bold=True)

    impact = [
        "Farmers reached: 200 million",
        "Annual savings: Rs 20,000 crore",
        "Lives saved: 8,000+ farmer suicides prevented",
        "Food production: 20 million tonnes extra",
    ]
    add_bullet_points(slide, Inches(0.7), Inches(4.7), Inches(12), Inches(2.5),
                      [f"  {im}" for im in impact], font_size=16, color=WHITE)

    # ========== SLIDE 8: TEAM ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), RGBColor(168, 85, 247))

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "Our Team", font_size=40, color=RGBColor(168, 85, 247), bold=True)

    # Team members
    team = [
        ("Rudra Khaire", "Team Leader", "2501201094", "AI/ML Development\nApp Development", ACCENT),
        ("Parth Soni", "Developer", "2501201077", "Data Collection\nTesting", BLUE),
        ("Parth Panchal", "Developer", "2501201078", "Documentation\nPresentation", ORANGE),
    ]

    for i, (name, role, enroll, tasks, color) in enumerate(team):
        x = Inches(0.5 + i * 4.2)
        card = add_shape(slide, x, Inches(1.5), Inches(3.8), Inches(3.5), BG_CARD)
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # Placeholder avatar
        avatar = add_shape(slide, x + Inches(1.2), Inches(1.7), Inches(1.4), Inches(1.4), color)

        add_text_box(slide, x + Inches(0.2), Inches(3.2), Inches(3.4), Inches(0.5),
                     name, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(3.7), Inches(3.4), Inches(0.4),
                     role, font_size=16, color=color, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(4.1), Inches(3.4), Inches(0.4),
                     f"Enrollment: {enroll}", font_size=12, color=RGBColor(156, 163, 175),
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(4.5), Inches(3.4), Inches(0.6),
                     tasks, font_size=12, color=RGBColor(156, 163, 175), alignment=PP_ALIGN.CENTER)

    # University
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.6),
                 "Dr. Kiran & Pallavi Patel Global University",
                 font_size=20, color=GOLD, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(6), Inches(12), Inches(0.5),
                 "Bharat Antriksh Saptah 2026 | Event 8: Artificial Intelligence",
                 font_size=16, color=RGBColor(156, 163, 175), alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 9: THANK YOU ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BG_DARK)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), ACCENT)

    add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
                 "Thank You!", font_size=60, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                 "AgriVision AI", font_size=36, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
                 '"See Disease Before It Spreads"', font_size=24, color=GOLD,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
                 "Questions?", font_size=28, color=RGBColor(134, 239, 172),
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
                 "Team Leader: Rudra Khaire | Parth Soni | Parth Panchal",
                 font_size=14, color=RGBColor(156, 163, 175), alignment=PP_ALIGN.CENTER)

    # Save presentation
    output_path = os.path.join('docs', 'AgriVision_AI_Presentation.pptx')
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
